from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor

def generate_presentation(
        ppt_heading, 
        explain_problem_you_are_solving, 
        target_market, your_product_service, 
        competetive_landscape, 
        market_validation, 
        revenue_model, 
        market_strategy, 
        team,
        financials,
        fund_requirement_deployment_plan):
    prs = Presentation()

    # --- Slide 1 ---
    blank_slide_layout = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide1.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/gov_logo.png", Inches(-0.25), Inches(-0.7), height=Inches(3))
    slide1.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/startupIndia.png", Inches(7.5), Inches(0.5), width=Inches(2))
    slide1.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(3.0), Inches(1), height=Inches(4))

    # --- Add Heading Text Below Center Image ---
    left = Inches(0.3)
    top = Inches(4.7)
    width = Inches(10)
    height = Inches(1)

    textbox = slide1.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "STARTUP INDIA SEED FUND SCHEME"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide1.shapes.add_textbox(left, Inches(5.5), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = ppt_heading
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    # --- Slide 2 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide2 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide2.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide2.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide2.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(0.1)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide2.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "EXPLAIN THE PROBLEM YOU ARE SOLVING"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide2.shapes.add_textbox(Inches(0.3), Inches(1.1), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = explain_problem_you_are_solving
    p.font.bold = False
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 3 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide3 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide3.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide3.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide3.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(0.1)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide3.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "TARGET MARKET"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide3.shapes.add_textbox(Inches(0.3), Inches(1.1), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = target_market
    p.font.bold = False
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 4 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide4 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide4.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide4.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide4.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(0.1)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide4.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "YOUR PRODUCT/SERVICE"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide4.shapes.add_textbox(Inches(0.3), Inches(1.1), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = your_product_service
    p.font.bold = False
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 5 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide5 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide5.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide5.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide5.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(0.1)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide5.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "COMPETITIVE LANDSCAPE"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide5.shapes.add_textbox(Inches(0.3), Inches(1.1), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = competetive_landscape
    p.font.bold = False
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 6 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide6 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide5.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(0.1)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide6.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "MARKET VALIDATION"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide6.shapes.add_textbox(Inches(0.3), Inches(1.1), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = market_validation
    p.font.bold = False
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 7 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide7 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide7.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide7.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide7.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(0.1)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide7.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "REVENUE MODEL"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide7.shapes.add_textbox(Inches(0.3), Inches(1.1), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = revenue_model
    p.font.bold = False
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 8 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide8 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide8.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide8.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide8.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(0.1)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide8.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "MARKET STRATEGY"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide8.shapes.add_textbox(Inches(0.3), Inches(1.1), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = market_strategy
    p.font.bold = False
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 9 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide9 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide9.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide9.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide9.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(0.1)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide9.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "TEAM"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide9.shapes.add_textbox(Inches(0.3), Inches(1.1), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = team
    p.font.bold = False
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 10 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide10 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide10.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide10.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide10.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(0.1)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide10.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "FINANCIALS"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide10.shapes.add_textbox(Inches(0.3), Inches(1.1), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = financials
    p.font.bold = False
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 11 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide11 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide11.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide11.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide11.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(0.1)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide11.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "FUND REQUIREMENT & DEPLOYMENT PLAN"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide11.shapes.add_textbox(Inches(0.3), Inches(1.1), width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = fund_requirement_deployment_plan
    p.font.bold = False
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 12 ---
    blank_slide_layout = prs.slide_layouts[6]
    slide12 = prs.slides.add_slide(blank_slide_layout)

    # --- Add Images ---
    slide12.shapes.add_picture("/Users/aditya.narayan/Desktop/form-to-ppt/asset/slide1/seedFundScheme.png", Inches(8.5), Inches(0.1), width=Inches(1.2))
    # Add a yellow line
    h_line = slide12.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.1), Inches(0.2), Inches(7), Inches(0.2))
    h_line.line.color.rgb = RGBColor(250, 191, 62)
    h_line.line.width = Pt(1)

    v_line = slide12.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.2), Inches(0.1), Inches(0.1), Inches(7))
    v_line.line.color.rgb = RGBColor(250, 191, 62)
    v_line.line.width = Pt(1)

    # --- Add Heading Text Below Center Image ---
    left = Inches(-0.5)
    top = Inches(3.5)
    width = Inches(9.5)
    height = Inches(1)

    textbox = slide12.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = "THANK YOU"
    p.font.bold = True
    p.font.size = Pt(35)
    p.font.color.rgb = RGBColor(47, 50, 144)
    p.alignment = PP_ALIGN.CENTER


    # Save the presentation
    prs.save("/Users/aditya.narayan/Desktop/form-to-ppt/output/auto-generated-ppt.pptx")
    print(f"🟢 Presentation saved as : auto-generated-ppt.pptx")